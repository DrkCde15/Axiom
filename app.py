from flask import Flask, render_template, request, jsonify
from google import genai
from google.genai import types
import os
import io

# Bibliotecas para extração de texto
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

ALLOWED_EXTENSIONS = {
    'png', 'jpg', 'jpeg', 'gif', 'webp', 
    'pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'txt'
}

MIME_TYPES = {
    'pdf': 'application/pdf',
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
    'txt': 'text/plain'
}

def extract_text_from_file(file, ext):
    try:
        if ext == 'txt':
            return file.read().decode('utf-8')
        elif ext in ['docx', 'doc']:
            doc = Document(file)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        elif ext in ['xlsx', 'xls']:
            df = pd.read_excel(file)
            return df.to_string()
        return ""
    except Exception as e:
        return f"[Erro ao extrair texto de {ext}: {str(e)}]"

GEMINI_API_KEY = os.getenv('GEMINI_API_KEY', 'SUA_API_KEY_AQUI')
client = genai.Client(api_key=GEMINI_API_KEY)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/chat')
def chat():
    return render_template('chat.html')

@app.route('/analyze', methods=['POST'])
def analyze():
    try:
        message = request.form.get('message')
        files = request.files.getlist('files')
        
        if not message:
            return jsonify({'error': 'Mensagem não fornecida'}), 400

        prompt = f"""Você é um assistente educacional especializado em ajudar estudantes com trabalhos escolares e universitários.

                    O aluno enviou a seguinte mensagem:
                    {message}

                    Instruções:
                    - Se for uma saudação (oi, olá, etc), responda de forma amigável e se apresente
                    - Se for uma pergunta, responda de forma clara e educativa
                    - Se for um arquivo, analise o arquivo e forneça feedback detalhado incluindo:
                    * Avaliação geral
                    * Pontos positivos
                    * Pontos a melhorar
                    * Sugestões específicas
                    * Nota de 0 a 10 e classificação (Excelente/Bom/Médio/Ruim)
                    - Se o aluno pedir para você fazer ou resolver algo, faça da melhor forma possível
                    - Se for uma conversa casual sobre estudos, seja prestativo e motivador

                    Responda de forma natural, clara e educativa."""
        contents = [prompt]
        
        for file in files:
            if file and file.filename:
                ext = file.filename.rsplit('.', 1)[1].lower()
                
                # Se for imagem ou PDF (suportados nativamente pelo Gemini)
                if ext in ['png', 'jpg', 'jpeg', 'gif', 'webp', 'pdf']:
                    file_data = file.read()
                    mime_type = MIME_TYPES.get(ext, 'application/octet-stream')
                    contents.append(types.Part.from_bytes(data=file_data, mime_type=mime_type))
                
                # Se for Office ou TXT, extraímos o texto e mandamos como string
                elif ext in ['docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls', 'txt']:
                    extracted_text = extract_text_from_file(file, ext)
                    contents.append(f"\n--- Conteúdo do arquivo {file.filename} ---\n{extracted_text}\n--- Fim do arquivo ---")

        response = client.models.generate_content(
            model='gemini-2.5-flash', # Nome estável do modelo
            contents=contents,
            config=types.GenerateContentConfig(temperature=0.7)
        )
        
        return jsonify({'success': True, 'response': response.text})
        
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)